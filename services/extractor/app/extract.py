# services/extractor/app/extract.py
import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import pdf2image
import pytesseract
from PIL import Image
import io
import nltk
from nltk.tokenize import sent_tokenize
from typing import List, Dict, Any
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.supported_extensions = {
            '.pdf', '.docx', '.pptx', '.txt', 
            '.jpg', '.jpeg', '.png',
        }
    
    def find_supported_files(self, directory: str) -> List[Path]:
        """Find all supported files in directory"""
        directory_path = Path(directory)
        
        if not directory_path.exists():
            raise FileNotFoundError(f"Directory not found: {directory}")
        
        files = []
        for ext in self.supported_extensions:
            # Search for both lowercase and uppercase extensions
            files.extend(directory_path.glob(f"*{ext}"))
            files.extend(directory_path.glob(f"*{ext.upper()}"))
            # Also search in subdirectories
            files.extend(directory_path.glob(f"**/*{ext}"))
            files.extend(directory_path.glob(f"**/*{ext.upper()}"))
        
        # Remove duplicates and return
        return list(set(files))
    
    async def process_file(self, file_path: str, chunk_size: int = 1000) -> List[Dict[str, Any]]:
        """Process a file and return text chunks with metadata"""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Extract text
        text = await self.extract_text_from_file(path)
        
        if not text.strip():
            logger.warning(f"No text extracted from {file_path}")
            return []
        
        # Chunk text
        chunks = self._chunk_text(text, chunk_size)
        
        return [
            {
                "text": chunk,
                "metadata": {
                    "filename": path.name,
                    "file_path": str(path),
                    "file_size": path.stat().st_size,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "file_type": path.suffix.lower()
                }
            }
            for i, chunk in enumerate(chunks)
        ]
    
    async def extract_text_from_file(self, file_path: Path) -> str:
        """Extract text from file path"""
        suffix = file_path.suffix.lower()
        
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
            
            if suffix == '.pdf':
                return self._extract_from_pdf(content)
            elif suffix == '.docx':
                return self._extract_from_docx(content)
            elif suffix == '.pptx':
                return self._extract_from_pptx(content)
            elif suffix == '.txt':
                return content.decode('utf-8', errors='ignore')
            elif suffix in ('.jpg', '.jpeg', '.png'):
                return self._extract_from_image(content)
            else:
                raise ValueError(f"Unsupported file format: {suffix}")
                
        except Exception as e:
            logger.error(f"Error extracting from {file_path}: {str(e)}")
            raise
    
    def _extract_from_pdf(self, content: bytes) -> str:
        """Extract text from PDF with OCR fallback"""
        try:
            # Method 1: Direct text extraction using PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            
            # If no text found or very little text, try OCR
            if not text.strip() or len(text.strip()) < 50:
                ocr_text = self._ocr_pdf(content)
                if ocr_text.strip():
                    text = ocr_text
                
            return text
            
        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            # Try OCR as fallback
            try:
                return self._ocr_pdf(content)
            except Exception as ocr_error:
                logger.error(f"PDF OCR also failed: {ocr_error}")
                raise e
    
    def _ocr_pdf(self, content: bytes) -> str:
        """OCR for PDFs without embedded text"""
        try:
            images = pdf2image.convert_from_bytes(content)
            text = ""
            for i, image in enumerate(images):
                logger.info(f"OCR processing page {i+1}/{len(images)}")
                page_text = pytesseract.image_to_string(image)
                text += f"Page {i+1}:\n{page_text}\n\n"
            return text
        except Exception as e:
            logger.error(f"PDF OCR failed: {str(e)}")
            return ""
    
    def _extract_from_docx(self, content: bytes) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(io.BytesIO(content))
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text += cell.text + "\n"
            
            return text
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            raise
    
    def _extract_from_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint"""
        try:
            prs = Presentation(io.BytesIO(content))
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"Slide {slide_num + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
            
            return text
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise
    
    def _extract_from_image(self, content: bytes) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"Image OCR failed: {str(e)}")
            raise
    
    def _chunk_text(self, text: str, chunk_size: int, chunk_overlap: int = 200) -> List[str]:
        """Split text into chunks with sentence awareness"""
        if not text.strip():
            return []
        
        # First, split into paragraphs
        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        
        chunks = []
        current_chunk = ""
        
        for paragraph in paragraphs:
            # If paragraph is too long, split into sentences
            if len(paragraph) > chunk_size:
                sentences = sent_tokenize(paragraph)
                for sentence in sentences:
                    if len(current_chunk) + len(sentence) <= chunk_size:
                        current_chunk += " " + sentence if current_chunk else sentence
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = sentence
            else:
                # Check if adding this paragraph would exceed chunk size
                if len(current_chunk) + len(paragraph) <= chunk_size:
                    current_chunk += " " + paragraph if current_chunk else paragraph
                else:
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                    current_chunk = paragraph
        
        # Don't forget the last chunk
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        # Ensure no chunk is too large
        final_chunks = []
        for chunk in chunks:
            if len(chunk) > chunk_size:
                # Split large chunks by sentences
                sentences = sent_tokenize(chunk)
                temp_chunk = ""
                for sentence in sentences:
                    if len(temp_chunk) + len(sentence) <= chunk_size:
                        temp_chunk += " " + sentence if temp_chunk else sentence
                    else:
                        if temp_chunk:
                            final_chunks.append(temp_chunk.strip())
                        temp_chunk = sentence
                if temp_chunk:
                    final_chunks.append(temp_chunk.strip())
            else:
                final_chunks.append(chunk)
        
        return final_chunks

    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return list(self.supported_extensions)